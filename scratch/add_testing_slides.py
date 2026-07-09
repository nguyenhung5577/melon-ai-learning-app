"""
Script thêm 3 slides Testing (cho Lâm Hoàng Vũ thuyết trình) vào file
[Presentation] Melon.pptx, chèn vào sau Slide 16 (Latency & Safety).
Slide mới:
  16.5 -> Testing Strategy Overview
  16.6 -> Test Metrics & Results
  16.7 -> Bugs Found & Fixed
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ─── COLORS (giữ nguyên brand Melon) ─────────────────
C_BG_DARK = RGBColor(0x0E, 0x0E, 0x1A)
C_BG_CARD = RGBColor(0x1C, 0x1C, 0x30)
C_MELON_G = RGBColor(0x4A, 0xDE, 0x80)
C_MELON_Y = RGBColor(0xFF, 0xD6, 0x00)
C_MELON_O = RGBColor(0xFF, 0x7A, 0x00)
C_MELON_P = RGBColor(0xA7, 0x8B, 0xFA)
C_MELON_B = RGBColor(0x38, 0xBD, 0xF8)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0x94, 0xA3, 0xB8)
C_BLACK   = RGBColor(0x0A, 0x0A, 0x0A)
C_RED     = RGBColor(0xFF, 0x45, 0x45)
C_ORANGE  = RGBColor(0xFF, 0x9A, 0x00)

W = Inches(13.33); H = Inches(7.5)

# ─── HELPERS ──────────────────────────────────────────
def slide_bg(slide, color=C_BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, border_color=None, border_w=Pt(2)):
    shape = slide.shapes.add_shape(1, l, t, w, h)
    sf = shape.fill
    if fill_color:
        sf.solid(); sf.fore_color.rgb = fill_color
    else:
        sf.background()
    line = shape.line
    if border_color:
        line.color.rgb = border_color; line.width = border_w
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(20), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run  = para.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = font_name
    return txb

def add_label_badge(slide, text, l, t, w=Inches(2.5), fill=C_MELON_B, text_color=C_BLACK):
    h = Inches(0.35)
    add_rect(slide, l, t, w, h, fill_color=fill, border_color=C_BLACK, border_w=Pt(1.5))
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    para = tf.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = text
    run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = text_color

def add_slide_at_end(prs, blank_layout):
    """Add new blank slide at the end."""
    return prs.slides.add_slide(blank_layout)

def insert_slide_after(prs, after_index):
    """
    Duplicate blank slide XML and insert it at a specific position.
    python-pptx doesn't support inserting at arbitrary positions natively,
    so we add at end, then move the XML element.
    """
    from pptx.oxml.ns import qn
    from lxml import etree

    # Add a new blank slide at end
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)

    # Get the slides XML list and reorder
    slides_xml = prs.slides._sldIdLst
    # The new slide is now the LAST element; move it to after_index+1
    items = list(slides_xml)
    last = items[-1]  # newly added
    slides_xml.remove(last)
    # Insert after `after_index` (0-based)
    ref_item = items[after_index]
    ref_item.addnext(last)

    return new_slide

# ─── LOAD EXISTING PRESENTATION ──────────────────────
src = r'd:\HCMUS\HK2-N3\SE4AI\melon-ai-learning-app\[Presentation] Melon.pptx'
prs = Presentation(src)
blank_layout = prs.slide_layouts[6]
print(f"Original slides: {len(prs.slides)}")

# ══════════════════════════════════════════════════════
# SLIDE A – TESTING STRATEGY OVERVIEW
# (Insert after slide 16, before Demo slides)
# ══════════════════════════════════════════════════════
s = insert_slide_after(prs, after_index=15)  # after slide index 15 = slide 16
slide_bg(s, C_BG_DARK)

add_label_badge(s, "05  TESTING STRATEGY", Inches(0.5), Inches(0.3),
                w=Inches(2.8), fill=C_MELON_B, text_color=C_BLACK)
add_text(s, "Chiến lược Kiểm thử Phần mềm", Inches(0.5), Inches(0.75),
         Inches(12), Inches(0.9), font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_B)

# Left: Test Approach
add_rect(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.1),
         fill_color=C_BG_CARD, border_color=C_MELON_B, border_w=Pt(2))
add_text(s, "🧪  Phương pháp Kiểm thử", Inches(0.7), Inches(2.0),
         Inches(5.4), Inches(0.55), font_size=Pt(16), bold=True, color=C_MELON_B)

approaches = [
    ("Manual Testing",   "Tester thực hiện thủ công trên\nChrome – localhost:3000"),
    ("Automation",       "Katalon Studio ghi & replay\n4 kịch bản tự động hóa"),
    ("API Testing",      "Postman test trực tiếp\nREST endpoints Backend"),
    ("Security Testing", "Thử nghiệm bypass Auth,\nIDOR, privilege escalation"),
]
for i, (title, desc) in enumerate(approaches):
    by = Inches(2.65) + i * Inches(1.05)
    col = [C_MELON_B, C_MELON_G, C_MELON_Y, C_MELON_O][i]
    add_rect(s, Inches(0.7), by, Inches(0.08), Inches(0.85), fill_color=col)
    add_text(s, title, Inches(0.9), by + Inches(0.02), Inches(4.8), Inches(0.4),
             font_size=Pt(13), bold=True, color=col)
    add_text(s, desc, Inches(0.9), by + Inches(0.45), Inches(4.8), Inches(0.5),
             font_size=Pt(11), color=C_GRAY)

# Right: Test Scope
add_rect(s, Inches(6.6), Inches(1.9), Inches(6.2), Inches(5.1),
         fill_color=C_BG_CARD, border_color=C_MELON_G, border_w=Pt(2))
add_text(s, "📋  Phạm vi Kiểm thử (5 Use-Cases)", Inches(6.8), Inches(2.0),
         Inches(5.8), Inches(0.55), font_size=Pt(14), bold=True, color=C_MELON_G)

scope_items = [
    ("UC1 – Authentication & Security",      "9 test cases",  C_MELON_G),
    ("UC2 – Features & Payments (Stripe)",   "6 test cases",  C_MELON_Y),
    ("UC3 – Knowledge Base & PDF Parsing",   "5 test cases",  C_MELON_O),
    ("UC4 – Practice & AI Personalization",  "9 test cases",  C_MELON_P),
    ("UC5 – Gamification & Admin System",    "19 test cases", C_MELON_B),
]
for i, (uc, cnt, col) in enumerate(scope_items):
    by = Inches(2.65) + i * Inches(0.85)
    add_rect(s, Inches(6.8), by, Inches(5.8), Inches(0.7),
             fill_color=C_BG_DARK, border_color=col, border_w=Pt(1.5))
    add_text(s, uc, Inches(6.95), by + Inches(0.1), Inches(4.2), Inches(0.5),
             font_size=Pt(12), color=C_WHITE)
    add_text(s, cnt, Inches(11.0), by + Inches(0.1), Inches(1.4), Inches(0.5),
             font_size=Pt(12), bold=True, color=col, align=PP_ALIGN.RIGHT)

# Bottom note
add_rect(s, Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.25),
         fill_color=C_BG_CARD, border_color=C_MELON_B, border_w=Pt(1))
add_text(s, "Tester: Lâm Hoàng Vũ  ·  Môi trường: Localhost (Next.js + FastAPI)  ·  Ngày: 21/06/2026",
         Inches(0.7), Inches(7.12), Inches(12), Inches(0.22),
         font_size=Pt(10), color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════
# SLIDE B – TEST METRICS & RESULTS
# ══════════════════════════════════════════════════════
s2 = insert_slide_after(prs, after_index=16)  # after the slide we just added (now index 16)
slide_bg(s2, C_BG_DARK)

add_label_badge(s2, "05  TEST RESULTS", Inches(0.5), Inches(0.3),
                w=Inches(2.8), fill=C_MELON_B, text_color=C_BLACK)
add_text(s2, "Kết quả Kiểm thử – Số liệu Tổng hợp", Inches(0.5), Inches(0.75),
         Inches(12), Inches(0.9), font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s2, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_B)

# Big KPIs row
kpis = [
    ("48", "Total\nTest Cases", C_MELON_B),
    ("27", "Passed ✅", C_MELON_G),
    ("8",  "Failed ❌", C_RED),
    ("13", "Skipped\n(Not Impl.)", C_GRAY),
    ("56%", "Pass Rate", C_MELON_Y),
]
for i, (num, label, col) in enumerate(kpis):
    bx = Inches(0.5) + i * Inches(2.55); by = Inches(1.9)
    add_rect(s2, bx, by, Inches(2.3), Inches(2.0),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_text(s2, num, bx, by + Inches(0.2), Inches(2.3), Inches(1.1),
             font_size=Pt(40), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s2, label, bx, by + Inches(1.3), Inches(2.3), Inches(0.6),
             font_size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER)

# Per Use-Case breakdown table
add_text(s2, "Breakdown theo Use-Case:", Inches(0.5), Inches(4.1),
         Inches(8), Inches(0.45), font_size=Pt(14), bold=True, color=C_MELON_B)

uc_data = [
    ("UC1 – Authentication & Security",    9,  8, 1, 0, C_MELON_G),
    ("UC2 – Features & Payments",          6,  4, 2, 0, C_MELON_Y),
    ("UC3 – Knowledge Base & Parsing",     5,  4, 1, 0, C_MELON_O),
    ("UC4 – Practice & Personalization",   9,  4, 1, 4, C_MELON_P),
    ("UC5 – Gamification & Admin",        19,  7, 3, 9, C_MELON_B),
]
hdrs = ["Use-Case", "Total", "✅ Pass", "❌ Fail", "⏭️ Skip"]
hw   = [4.5, 1.2, 1.4, 1.4, 1.4]
hx = 0.5
for hdr, hw_ in zip(hdrs, hw):
    add_rect(s2, Inches(hx), Inches(4.6), Inches(hw_ - 0.05), Inches(0.4),
             fill_color=C_MELON_B, border_color=C_BLACK, border_w=Pt(1))
    add_text(s2, hdr, Inches(hx), Inches(4.6), Inches(hw_ - 0.05), Inches(0.4),
             font_size=Pt(11), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    hx += hw_

for i, (name, total, passed, failed, skip, col) in enumerate(uc_data):
    by = Inches(5.06) + i * Inches(0.42)
    row_data = [name, str(total), str(passed), str(failed), str(skip)]
    rx = 0.5
    for j, (val, hw_) in enumerate(zip(row_data, hw)):
        fc = C_BG_CARD if j > 0 else C_BG_DARK
        if j == 3 and int(val) > 0: fc = RGBColor(0x2B, 0x08, 0x08)
        add_rect(s2, Inches(rx), by, Inches(hw_ - 0.05), Inches(0.38),
                 fill_color=fc, border_color=col if j == 0 else C_BG_CARD, border_w=Pt(1.5))
        tc = col if j == 0 else (C_MELON_G if j == 2 else (C_RED if j == 3 and int(val) > 0 else C_GRAY))
        add_text(s2, val, Inches(rx), by, Inches(hw_ - 0.05), Inches(0.38),
                 font_size=Pt(11), color=tc, align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT,
                 bold=(j == 0))
        rx += hw_

# ══════════════════════════════════════════════════════
# SLIDE C – BUGS FOUND & FIXED
# ══════════════════════════════════════════════════════
s3 = insert_slide_after(prs, after_index=17)
slide_bg(s3, C_BG_DARK)

add_label_badge(s3, "05  BUG REPORT & FIXES", Inches(0.5), Inches(0.3),
                w=Inches(3.0), fill=C_MELON_B, text_color=C_BLACK)
add_text(s3, "8 Bugs Phát hiện  ·  Đã Fix Critical", Inches(0.5), Inches(0.75),
         Inches(12), Inches(0.9), font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s3, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_B)

# Summary severity badges
severities = [
    ("🚨  2 CRITICAL", C_RED),
    ("🔴  1 HIGH",     C_MELON_O),
    ("🟡  3 MEDIUM",   C_MELON_Y),
    ("🟢  2 LOW",      C_MELON_G),
]
for i, (label, col) in enumerate(severities):
    bx = Inches(0.5) + i * Inches(3.1)
    add_rect(s3, bx, Inches(1.9), Inches(2.9), Inches(0.6),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    add_text(s3, label, bx, Inches(1.9), Inches(2.9), Inches(0.6),
             font_size=Pt(14), bold=True, color=col, align=PP_ALIGN.CENTER)

# Bug list
bugs = [
    ("SEC-01", "IDOR – Xem trộm điểm học sinh",
     "Không verify UID trên API progress", "CRITICAL", "✅ FIXED",
     "Thêm requireChildAccess middleware", C_RED),
    ("SEC-02", "Backdoor tự thăng cấp Admin",
     "Dev Helper trên /admin cho set role=admin", "CRITICAL", "✅ FIXED",
     "Xóa Dev Helper, chặn Firestore rules", C_RED),
    ("ERR-03", "Hủy Pro bị cắt quyền ngay lập tức",
     "Webhook không check cancel_at_period_end", "MEDIUM", "✅ FIXED",
     "Giữ Pro đến hết kỳ thanh toán", C_MELON_Y),
    ("ERR-03b","Pro không đồng bộ Admin dashboard",
     "Checkout chỉ update collection subscriptions", "HIGH", "✅ FIXED",
     "Mirror isPro sang collection users", C_MELON_O),
    ("ERR-04", "File >20MB upload không bị chặn",
     "Frontend & API thiếu kiểm tra file.size", "MEDIUM", "✅ FIXED",
     "Chặn ngay ở UI + API trả 413", C_MELON_Y),
    ("UX-01",  "Thiếu nút TTS đọc gợi ý AI",
     "PersonalizedExercisePanel quên gắn nút Loa", "LOW", "✅ FIXED",
     "Thêm Speaker button + gọi /api/v1/ai/tts", C_MELON_G),
    ("FAM-01", "State Sync mâu thuẫn số con",
     "Dùng nhiều biến state khác nhau cho cùng 1 list", "MEDIUM", "✅ FIXED",
     "Đồng bộ về 1 biến children duy nhất", C_MELON_Y),
    ("UX-02",  "Sai tiêu đề Profile Phụ huynh",
     "Hardcode chữ \"Hồ sơ của con\" cho mọi role", "LOW", "✅ FIXED",
     "Render tiêu đề theo role (kid/parent/admin)", C_MELON_G),
]

# Show 8 bugs in compact format
for i, (code, title, cause, sev, status, fix, col) in enumerate(bugs):
    row = i // 2; ci = i % 2
    bx = Inches(0.5) + ci * Inches(6.2); by = Inches(2.65) + row * Inches(1.15)
    bw = Inches(6.0); bh = Inches(1.05)
    add_rect(s3, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(1.5))
    # Code badge
    add_rect(s3, bx, by, Inches(1.0), bh, fill_color=col)
    add_text(s3, code, bx + Inches(0.05), by + Inches(0.3), Inches(0.9), Inches(0.45),
             font_size=Pt(10), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    # Title
    add_text(s3, title, bx + Inches(1.1), by + Inches(0.05), bw - Inches(2.2), Inches(0.45),
             font_size=Pt(12), bold=True, color=col)
    # Fix
    add_text(s3, "→ " + fix, bx + Inches(1.1), by + Inches(0.55), bw - Inches(2.2), Inches(0.45),
             font_size=Pt(10), color=C_GRAY)
    # Status
    add_text(s3, status, bx + bw - Inches(1.3), by + Inches(0.28), Inches(1.2), Inches(0.45),
             font_size=Pt(12), bold=True, color=C_MELON_G, align=PP_ALIGN.RIGHT)

# ─── SAVE ─────────────────────────────────────────────
out_path = r'd:\HCMUS\HK2-N3\SE4AI\melon-ai-learning-app\[Presentation] Melon.pptx'
prs.save(out_path)
print(f"✅  Saved: {out_path}")
print(f"   Total slides now: {len(prs.slides)}")
print("   New slides added: Testing Strategy, Test Metrics, Bug Report")
