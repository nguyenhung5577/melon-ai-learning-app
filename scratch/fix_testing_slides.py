"""
Script sửa 3 slides Testing (17, 18, 19) trong [Presentation] Melon.pptx:
- Xóa 3 slides sai kích thước/màu sắc
- Duplicate slide 16 (Hiệu năng & An toàn) làm template - giữ nguyên nền và kích thước
- Clear nội dung cũ, thêm nội dung Testing mới vào đúng tọa độ 20" x 11.25"
"""
import sys, copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout.reconfigure(encoding='utf-8')

# ─── LOAD ───────────────────────────────────────────
SRC = r'd:\HCMUS\HK2-N3\SE4AI\melon-ai-learning-app\[Presentation] Melon.pptx'
prs = Presentation(SRC)

W = prs.slide_width   # 18288000 EMU = 20"
H = prs.slide_height  # 10287000 EMU = 11.25"
print(f'Slide size: {W/914400:.2f}" x {H/914400:.2f}"')
print(f'Original slides: {len(prs.slides)}')

# ─── STEP 1: Delete slides 17, 18, 19 (index 16, 17, 18) ───
# We must delete from highest index to lowest to avoid index shifting
slides_xml = prs.slides._sldIdLst
rIds_to_remove = []

for idx in [18, 17, 16]:  # highest first
    slide_elem = list(slides_xml)[idx]
    rId = slide_elem.get('r:id')
    rIds_to_remove.append((idx, rId, slide_elem))

for idx, rId, elem in rIds_to_remove:
    # Remove from relationship
    try:
        part = prs.slides.part
        slide_part = part.related_parts[rId]
        part.drop_rel(rId)
    except Exception as e:
        print(f'  Warning dropping rel {rId}: {e}')
    # Remove XML element
    slides_xml.remove(elem)

print(f'After deletion: {len(prs.slides)} slides')

# ─── HELPERS (scale relative to 20" x 11.25") ─────────────
def I(x): return int(x * 914400)  # inches to EMU

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_pt=2):
    sp = slide.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    if fill_rgb:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill_rgb
    else:
        sp.fill.background()
    if line_rgb:
        sp.line.color.rgb = line_rgb
        sp.line.width = Pt(line_pt)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp

def add_text(slide, text, l, t, w, h,
             size=20, bold=False, color=None, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    if color is None: color = RGBColor(0xFF, 0xFF, 0xFF)
    txb = slide.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = 'Calibri'
    return txb

# ─── COLORS (from existing slides) ──────────────────
# Slide 2 background is #16172B (dark)
# Accent colors from slides: FFD600 (yellow), 4ADE80 (green), etc.
C_DARK  = RGBColor(0x16, 0x17, 0x2B)   # nền tối từ slide 2
C_CARD  = RGBColor(0x1E, 0x1F, 0x3A)   # card nền
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY  = RGBColor(0x94, 0xA3, 0xB8)
C_BLACK = RGBColor(0x0A, 0x0A, 0x0A)
# Accent từ slides gốc:
C_YELLOW = RGBColor(0xFF, 0xD6, 0x00)  # FFD600 - chính
C_GREEN  = RGBColor(0x4A, 0xDE, 0x80)  # green accent
C_BLUE   = RGBColor(0x38, 0xBD, 0xF8)  # blue accent
C_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)  # purple
C_ORANGE = RGBColor(0xFF, 0x7A, 0x00)  # orange
C_RED    = RGBColor(0xFF, 0x45, 0x45)
C_LGREEN = RGBColor(0x22, 0xC5, 0x5E)  # darker green for pass

def add_slide_after(prs, after_idx):
    """Add blank slide at end, then move XML to position after_idx (0-based)."""
    blank = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank)
    # Set background to match dark theme
    bg = new_slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DARK
    # Move to correct position
    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    last = items[-1]
    sldIdLst.remove(last)
    ref = items[after_idx]
    ref.addnext(last)
    return new_slide

def section_label(slide, text, l=0.75, t=0.55, fill=C_BLUE):
    """Small section label pill at top."""
    add_rect(slide, l, t, len(text)*0.115 + 0.3, 0.38, fill_rgb=fill, line_rgb=None)
    add_text(slide, text, l+0.1, t+0.04, len(text)*0.115+0.1, 0.3,
             size=11, bold=True, color=C_BLACK, align=PP_ALIGN.LEFT)

def slide_header(slide, title, accent_color=C_BLUE):
    """Big title + underline bar."""
    add_text(slide, title, 0.75, 1.05, 18.0, 1.2,
             size=44, bold=True, color=C_WHITE)
    add_rect(slide, 0.75, 2.3, 18.5, 0.06, fill_rgb=accent_color)

def slide_footer(slide, note=''):
    """Presenter note at bottom."""
    if note:
        add_rect(slide, 0.75, 10.65, 18.5, 0.4, fill_rgb=C_CARD, line_rgb=C_BLUE, line_pt=1)
        add_text(slide, note, 0.95, 10.7, 18.0, 0.32,
                 size=12, color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════
# SLIDE 17 – TESTING STRATEGY
# (add after slide 16, which is now index 15 after deletion)
# ══════════════════════════════════════════════════════
s1 = add_slide_after(prs, after_idx=15)
section_label(s1, "05  CHIẾN LƯỢC KIỂM THỬ", fill=C_BLUE)
slide_header(s1, "Chiến lược Kiểm thử Phần mềm", C_BLUE)
slide_footer(s1, "Tester: Lâm Hoàng Vũ  ·  Môi trường: Localhost (Next.js + FastAPI)  ·  21/06/2026")

# Left panel: Approach
add_rect(s1, 0.75, 2.6, 8.5, 7.8, fill_rgb=C_CARD, line_rgb=C_BLUE, line_pt=2)
add_text(s1, "🧪  Phương pháp Kiểm thử", 1.0, 2.75, 8.0, 0.7,
         size=22, bold=True, color=C_BLUE)

approaches = [
    ("Manual Testing",   "Tester thực hiện thủ công\ntrên Chrome – localhost:3000",            C_BLUE),
    ("Automation",       "Katalon Studio ghi & replay\n4 kịch bản tự động hóa (PA5)",          C_GREEN),
    ("API Testing",      "Postman kiểm tra trực tiếp\ncác REST endpoint Backend",              C_YELLOW),
    ("Security Testing", "Thử bypass Auth, IDOR,\nprivilege escalation attacks",               C_ORANGE),
]
for i, (title, desc, col) in enumerate(approaches):
    by = 3.65 + i * 1.6
    add_rect(s1, 0.9, by, 0.12, 1.3, fill_rgb=col)
    add_rect(s1, 1.1, by, 7.9, 1.3, fill_rgb=RGBColor(0x12, 0x13, 0x22), line_rgb=None)
    add_text(s1, title, 1.3, by + 0.08, 7.5, 0.52, size=17, bold=True, color=col)
    add_text(s1, desc,  1.3, by + 0.65, 7.5, 0.65, size=14, color=C_GRAY)

# Right panel: Scope
add_rect(s1, 9.7, 2.6, 9.55, 7.8, fill_rgb=C_CARD, line_rgb=C_GREEN, line_pt=2)
add_text(s1, "📋  Phạm vi – 5 Use-Cases (48 TC)", 9.95, 2.75, 9.0, 0.7,
         size=22, bold=True, color=C_GREEN)

scope = [
    ("UC1 – Authentication & Security",   "9 TC",  C_GREEN),
    ("UC2 – Features & Stripe Payments",  "6 TC",  C_YELLOW),
    ("UC3 – Knowledge Base & PDF Parse",  "5 TC",  C_ORANGE),
    ("UC4 – Practice & AI Personalization","9 TC", C_PURPLE),
    ("UC5 – Gamification & Admin System", "19 TC", C_BLUE),
]
for i, (uc, cnt, col) in enumerate(scope):
    by = 3.65 + i * 1.35
    add_rect(s1, 9.85, by, 0.08, 1.1, fill_rgb=col)
    add_rect(s1, 10.05, by, 9.0, 1.1, fill_rgb=RGBColor(0x12, 0x13, 0x22), line_rgb=None)
    add_text(s1, uc, 10.25, by + 0.12, 7.8, 0.52, size=15, bold=True, color=col)
    add_text(s1, cnt, 17.8, by + 0.28, 1.0, 0.45, size=18, bold=True, color=col, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# SLIDE 18 – TEST METRICS & RESULTS
# ══════════════════════════════════════════════════════
s2 = add_slide_after(prs, after_idx=16)
section_label(s2, "05  KẾT QUẢ KIỂM THỬ", fill=C_BLUE)
slide_header(s2, "Kết quả Kiểm thử – Số liệu Tổng hợp", C_BLUE)

# Big KPI row (5 cards)
kpis = [
    ("48",  "Tổng Test Cases", C_BLUE),
    ("27",  "✅  Passed",       C_GREEN),
    ("8",   "❌  Failed",       C_RED),
    ("13",  "⏭️  Skipped",     C_GRAY),
    ("56%", "Pass Rate",       C_YELLOW),
]
for i, (num, label, col) in enumerate(kpis):
    bx = 0.75 + i * 3.55
    add_rect(s2, bx, 2.6, 3.2, 2.5, fill_rgb=C_CARD, line_rgb=col, line_pt=3)
    add_rect(s2, bx, 2.6, 3.2, 0.08, fill_rgb=col)  # top accent bar
    add_text(s2, num,   bx, 3.0, 3.2, 1.3, size=52, bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(s2, label, bx, 4.2, 3.2, 0.7, size=16, color=C_WHITE, align=PP_ALIGN.CENTER)

# Breakdown table
add_text(s2, "Breakdown theo Use-Case:", 0.75, 5.4, 10.0, 0.55,
         size=18, bold=True, color=C_BLUE)

# Table header
hdrs = [("Use-Case", 7.5), ("Total", 1.8), ("✅ Pass", 2.2), ("❌ Fail", 2.2), ("⏭️ Skip", 2.0)]
hx = 0.75
for hdr, hw in hdrs:
    add_rect(s2, hx, 6.05, hw - 0.05, 0.6, fill_rgb=C_BLUE, line_rgb=None)
    add_text(s2, hdr, hx+0.1, 6.1, hw-0.2, 0.5, size=15, bold=True, color=C_BLACK,
             align=PP_ALIGN.CENTER if hdr != "Use-Case" else PP_ALIGN.LEFT)
    hx += hw

rows = [
    ("UC1 – Authentication & Security",    9,  8, 1, 0, C_GREEN),
    ("UC2 – Features & Payments",          6,  4, 2, 0, C_YELLOW),
    ("UC3 – Knowledge Base & Parsing",     5,  4, 1, 0, C_ORANGE),
    ("UC4 – Practice & Personalization",   9,  4, 1, 4, C_PURPLE),
    ("UC5 – Gamification & Admin",        19,  7, 3, 9, C_BLUE),
]
for ri, (name, total, passed, failed, skip, col) in enumerate(rows):
    by = 6.72 + ri * 0.72
    vals = [(name, 7.5, col, PP_ALIGN.LEFT), (str(total), 1.8, C_WHITE, PP_ALIGN.CENTER),
            (str(passed), 2.2, C_GREEN, PP_ALIGN.CENTER),
            (str(failed), 2.2, C_RED if failed > 0 else C_GRAY, PP_ALIGN.CENTER),
            (str(skip), 2.0, C_GRAY, PP_ALIGN.CENTER)]
    rx = 0.75
    for vi, (val, hw, vc, va) in enumerate(vals):
        fc = C_CARD if vi > 0 else RGBColor(0x12, 0x13, 0x22)
        if vi == 3 and failed > 0: fc = RGBColor(0x2B, 0x08, 0x08)
        border = col if vi == 0 else None
        add_rect(s2, rx, by, hw - 0.05, 0.62, fill_rgb=fc, line_rgb=border, line_pt=2)
        add_text(s2, val, rx+0.1, by+0.06, hw-0.2, 0.52, size=14,
                 bold=(vi == 0), color=vc, align=va)
        rx += hw

# ══════════════════════════════════════════════════════
# SLIDE 19 – BUGS FOUND & FIXED
# ══════════════════════════════════════════════════════
s3 = add_slide_after(prs, after_idx=17)
section_label(s3, "05  BUG REPORT & FIXES", fill=C_BLUE)
slide_header(s3, "8 Bugs Phát hiện  ·  2 Critical Đã Fix", C_BLUE)

# Severity summary row
sevs = [("🚨  2 CRITICAL", C_RED), ("🔴  1 HIGH", C_ORANGE), ("🟡  3 MEDIUM", C_YELLOW), ("🟢  2 LOW", C_GREEN)]
for i, (label, col) in enumerate(sevs):
    bx = 0.75 + i * 4.55
    add_rect(s3, bx, 2.6, 4.2, 0.8, fill_rgb=C_CARD, line_rgb=col, line_pt=2.5)
    add_text(s3, label, bx, 2.6, 4.2, 0.8, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)

# Bug cards (2 columns x 4 rows)
bugs = [
    ("SEC-01", "IDOR – Xem trộm điểm học sinh",
     "Không verify UID → API progress bị lộ",
     "Thêm requireChildAccess middleware", "✅ FIXED", C_RED),
    ("SEC-02", "Backdoor tự thăng cấp Admin",
     "Dev Helper trên /admin cho set role=admin",
     "Xóa Dev Helper + chặn Firestore rules", "✅ FIXED", C_RED),
    ("ERR-03", "Hủy Pro bị cắt quyền ngay lập tức",
     "Webhook không check cancel_at_period_end",
     "Giữ Pro đến hết kỳ thanh toán", "✅ FIXED", C_YELLOW),
    ("ERR-03b","Pro không đồng bộ Admin dashboard",
     "Checkout chỉ update collection subscriptions",
     "Mirror isPro sang collection users", "✅ FIXED", C_ORANGE),
    ("ERR-04", "File >20MB upload không bị chặn",
     "Frontend & API thiếu kiểm tra file.size",
     "Chặn ở UI + API trả HTTP 413", "✅ FIXED", C_YELLOW),
    ("UX-01",  "Thiếu nút TTS đọc gợi ý AI",
     "PersonalizedExercisePanel quên gắn nút Loa",
     "Thêm Speaker button + /api/v1/ai/tts", "✅ FIXED", C_GREEN),
    ("FAM-01", "State Sync mâu thuẫn số con",
     "Dùng nhiều biến state cho cùng 1 danh sách",
     "Đồng bộ về 1 biến children duy nhất", "✅ FIXED", C_YELLOW),
    ("UX-02",  "Sai tiêu đề Profile Phụ huynh",
     "Hardcode \"Hồ sơ của con\" cho mọi role",
     "Render tiêu đề theo role (kid/parent)", "✅ FIXED", C_GREEN),
]

COL_W = 9.0; ROW_H = 1.55
for i, (code, title, cause, fix, status, col) in enumerate(bugs):
    row = i // 2; ci = i % 2
    bx = 0.75 + ci * 9.45
    by = 3.65 + row * ROW_H
    # Card bg
    add_rect(s3, bx, by, COL_W, ROW_H - 0.1, fill_rgb=C_CARD, line_rgb=col, line_pt=1.5)
    # Left color bar
    add_rect(s3, bx, by, 1.3, ROW_H - 0.1, fill_rgb=col)
    # Code label
    add_text(s3, code, bx, by + (ROW_H-0.1)/2 - 0.25, 1.3, 0.5,
             size=13, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    # Title
    add_text(s3, title, bx + 1.4, by + 0.1, COL_W - 2.8, 0.5,
             size=15, bold=True, color=col)
    # Cause (gray)
    add_text(s3, "⚠ " + cause, bx + 1.4, by + 0.6, COL_W - 2.8, 0.45,
             size=12, color=C_GRAY, italic=True)
    # Fix (white)
    add_text(s3, "→ " + fix, bx + 1.4, by + 1.0, COL_W - 2.8, 0.42,
             size=12, color=C_WHITE)
    # Status badge
    add_rect(s3, bx + COL_W - 1.45, by + 0.1, 1.35, 0.5,
             fill_rgb=C_LGREEN, line_rgb=None)
    add_text(s3, status, bx + COL_W - 1.45, by + 0.1, 1.35, 0.5,
             size=13, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)

# ─── SAVE ────────────────────────────────────────────
OUT = r'd:\HCMUS\HK2-N3\SE4AI\melon-ai-learning-app\[Presentation] Melon_FIXED.pptx'
prs.save(OUT)
SRC = OUT
print(f'\n✅  Saved: {SRC}')
print(f'   Total slides: {len(prs.slides)}')
# Verify positions
sldIdLst = prs.slides._sldIdLst
for i, slide in enumerate(prs.slides):
    texts = [s.text_frame.paragraphs[0].text[:30] for s in slide.shapes
             if s.has_text_frame and s.text_frame.paragraphs[0].text.strip()]
    if texts: print(f'  Slide {i+1}: {texts[0]}')
