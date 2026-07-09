"""
Script tạo slides thuyết trình Melon AI Learning App (20 slides) bằng python-pptx.
Thiết kế: Hiện đại, nhiều màu sắc, tập trung vào hình ảnh/biểu đồ minh họa.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy

# ─── COLORS ───────────────────────────────────────────
C_BG_DARK    = RGBColor(0x0E, 0x0E, 0x1A)   # Nền tối (xanh đêm)
C_BG_CARD    = RGBColor(0x1C, 0x1C, 0x30)   # Card nền
C_MELON_G    = RGBColor(0x4A, 0xDE, 0x80)   # Xanh lá Melon (neon green)
C_MELON_Y    = RGBColor(0xFF, 0xD6, 0x00)   # Vàng Melon
C_MELON_O    = RGBColor(0xFF, 0x7A, 0x00)   # Cam Melon
C_MELON_P    = RGBColor(0xA7, 0x8B, 0xFA)   # Tím Melon
C_MELON_B    = RGBColor(0x38, 0xBD, 0xF8)   # Xanh dương
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY       = RGBColor(0x94, 0xA3, 0xB8)
C_BORDER     = RGBColor(0xFF, 0xD6, 0x00)   # Border Neobrutalism
C_BLACK      = RGBColor(0x0A, 0x0A, 0x0A)

# ─── HELPERS ──────────────────────────────────────────
def slide_bg(slide, color=C_BG_DARK):
    """Fill slide background with solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color=None, border_color=None, border_w=Pt(2)):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    sf = shape.fill
    if fill_color:
        sf.solid(); sf.fore_color.rgb = fill_color
    else:
        sf.background()
    line = shape.line
    if border_color:
        line.color.rgb = border_color
        line.width = border_w
    else:
        line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=Pt(20), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run  = para.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name  = font_name
    return txb

def add_label_badge(slide, text, l, t, fill=C_MELON_G, text_color=C_BLACK):
    """Small pill/badge label."""
    w = Inches(1.8); h = Inches(0.35)
    r = add_rect(slide, l, t, w, h, fill_color=fill, border_color=C_BLACK, border_w=Pt(1.5))
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run(); run.text = text
    run.font.size = Pt(10); run.font.bold = True; run.font.color.rgb = text_color
    return r

# ─── LAYOUT ───────────────────────────────────────────
W = Inches(13.33); H = Inches(7.5)   # 16:9 widescreen

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]   # completely blank

# ══════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

# Big decorative circle top-right
c = s.shapes.add_shape(9, Inches(9.5), Inches(-1.2), Inches(5), Inches(5))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x2A)
c.line.fill.background()

# Melon logo text
add_text(s, "🍉 MELON", Inches(1), Inches(0.5), Inches(4), Inches(0.7),
         font_size=Pt(22), bold=True, color=C_MELON_G, font_name="Calibri")

# Main title
add_text(s, "AI Learning App", Inches(1), Inches(1.2), Inches(8), Inches(1.3),
         font_size=Pt(54), bold=True, color=C_WHITE, font_name="Calibri")

# Subtitle
add_text(s, "Nền tảng Giáo dục Tương tác ứng dụng AI cho Học sinh",
         Inches(1), Inches(2.6), Inches(9), Inches(0.6),
         font_size=Pt(20), color=C_MELON_Y, font_name="Calibri")

# Divider line
add_rect(s, Inches(1), Inches(3.3), Inches(5), Inches(0.04),
         fill_color=C_MELON_G)

# Team info
team_lines = [
    "Nhóm 2  |  SE4AI – HK2-2026",
    "Nguyễn Huy Hoàng (23122031)   Nguyễn Bá Nam (23122043)",
    "Trần Tạ Quang Minh (23122042)  Nguyễn Đăng Khôi (23122037)",
    "Lâm Hoàng Vũ (23122056)",
]
for i, line in enumerate(team_lines):
    add_text(s, line, Inches(1), Inches(3.55) + Inches(0.42)*i, Inches(10), Inches(0.42),
             font_size=Pt(13), color=C_GRAY if i > 0 else C_WHITE, font_name="Calibri")

# Bottom tag
add_label_badge(s, "10/7/2026 · Online", Inches(1), Inches(6.9), fill=C_MELON_P, text_color=C_WHITE)

# ══════════════════════════════════════════════════════
# SLIDE 2 – AGENDA
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_text(s, "📋  Nội dung Thuyết trình", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
         font_size=Pt(32), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.04), fill_color=C_MELON_G)

sections = [
    ("01", "Vision & Problem Statement",  "TM1 – Hoàng",  C_MELON_G),
    ("02", "Project Management & Use Cases","TM2 – Nam",   C_MELON_Y),
    ("03", "Architecture & NFRs",          "TM3 – Minh",  C_MELON_O),
    ("04", "Tech Stack & AI Integration",  "TM4 – Khôi",  C_MELON_P),
    ("05", "Evaluation & Experiments",     "TM5 – Vũ",    C_MELON_B),
    ("06", "Demo & Conclusion",            "TM6 – Hùng",  C_MELON_G),
]
cols = 3
for idx, (num, title, member, col) in enumerate(sections):
    row = idx // cols; col_i = idx % cols
    bx = Inches(0.5) + col_i * Inches(4.3); by = Inches(1.4) + row * Inches(2.4)
    bw = Inches(4.0); bh = Inches(2.1)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_text(s, num, bx + Inches(0.15), by + Inches(0.1), Inches(0.8), Inches(0.6),
             font_size=Pt(28), bold=True, color=col)
    add_text(s, title, bx + Inches(0.15), by + Inches(0.7), bw - Inches(0.3), Inches(0.8),
             font_size=Pt(14), bold=True, color=C_WHITE)
    add_text(s, member, bx + Inches(0.15), by + Inches(1.5), bw - Inches(0.3), Inches(0.45),
             font_size=Pt(11), color=C_GRAY, italic=True)

# ══════════════════════════════════════════════════════
# SLIDE 3 – PROBLEM STATEMENT
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "01  VISION", Inches(0.5), Inches(0.3), fill=C_MELON_G, text_color=C_BLACK)
add_text(s, "Bài toán cần giải quyết", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_G)

problems = [
    ("😴", "Trẻ nhàm chán",
     "Cách học đọc-chép truyền thống\nquá đơn điệu → mất tập trung,\nthiếu hứng thú.", C_MELON_O),
    ("🔒", "Thiếu an toàn",
     "Phụ huynh lo ngại nội dung\nkhông kiểm soát trên mạng,\nkhó giám sát trẻ.", C_MELON_P),
    ("📉", "Không cá nhân hóa",
     "Chương trình học đồng loạt,\nkhông phù hợp tốc độ &\nnăng lực từng trẻ.", C_MELON_Y),
]
for i, (icon, title, desc, col) in enumerate(problems):
    bx = Inches(0.5) + i * Inches(4.25); by = Inches(1.9)
    add_rect(s, bx, by, Inches(4.0), Inches(5.0),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_text(s, icon, bx + Inches(0.15), by + Inches(0.2), Inches(1), Inches(0.9),
             font_size=Pt(36))
    add_text(s, title, bx + Inches(0.15), by + Inches(1.15), Inches(3.7), Inches(0.6),
             font_size=Pt(17), bold=True, color=col)
    add_text(s, desc, bx + Inches(0.15), by + Inches(1.8), Inches(3.7), Inches(2.8),
             font_size=Pt(13), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 4 – PRODUCT POSITION STATEMENT
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "01  VISION", Inches(0.5), Inches(0.3), fill=C_MELON_G, text_color=C_BLACK)
add_text(s, "Product Position Statement", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_G)

# Big quote
add_rect(s, Inches(0.5), Inches(1.9), Inches(12.3), Inches(1.6),
         fill_color=RGBColor(0x10, 0x2A, 0x1A), border_color=C_MELON_G, border_w=Pt(2))
add_text(s, "\"Học thông minh  ·  Vui hơn  ·  Tiến bộ hơn\"",
         Inches(0.8), Inches(2.05), Inches(11.8), Inches(1.2),
         font_size=Pt(28), bold=True, color=C_MELON_G, align=PP_ALIGN.CENTER)

# Template table
rows = [
    ("Dành cho",   "Trẻ em 6-12 tuổi & Phụ huynh"),
    ("Đây là",     "Nền tảng giáo dục tương tác tích hợp GenAI"),
    ("Khác biệt",  "Biến mọi tài liệu PDF/Ảnh thành bài học tương tác, gamified"),
    ("Giá trị cốt lõi", "AI Tutor cá nhân + Child-Safety Moderation + Real-time Progress"),
]
for i, (label, val) in enumerate(rows):
    by = Inches(3.65) + i * Inches(0.75)
    add_rect(s, Inches(0.5), by, Inches(2.5), Inches(0.65),
             fill_color=C_MELON_G, border_color=C_BLACK, border_w=Pt(1))
    add_text(s, label, Inches(0.6), by + Inches(0.12), Inches(2.3), Inches(0.5),
             font_size=Pt(12), bold=True, color=C_BLACK)
    add_rect(s, Inches(3.1), by, Inches(9.7), Inches(0.65),
             fill_color=C_BG_CARD, border_color=C_MELON_G, border_w=Pt(1))
    add_text(s, val, Inches(3.25), by + Inches(0.12), Inches(9.4), Inches(0.5),
             font_size=Pt(13), color=C_WHITE)

# ══════════════════════════════════════════════════════
# SLIDE 5 – USERS & MARKET
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "01  VISION", Inches(0.5), Inches(0.3), fill=C_MELON_G, text_color=C_BLACK)
add_text(s, "Người dùng & Thị trường", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_G)

users = [
    ("🧒", "KID\n(Học sinh)", [
        "Trẻ 6–12 tuổi",
        "Học qua flashcard, kéo-thả",
        "Nghe TTS đọc đề bài",
        "Nhận XP, Badge, Leaderboard",
    ], C_MELON_G),
    ("👨‍👩‍👧", "PARENT\n(Phụ huynh)", [
        "Quản lý tài khoản con",
        "Upload PDF → AI sinh bài",
        "Xem Dashboard tiến độ",
        "Nhận thông báo học tập",
    ], C_MELON_Y),
    ("🛡️", "ADMIN\n(Quản trị)", [
        "Kiểm duyệt nội dung AI",
        "Quản lý Question Bank",
        "Upload & Parse đề thi",
        "Giám sát hệ thống",
    ], C_MELON_P),
]
for i, (icon, title, points, col) in enumerate(users):
    bx = Inches(0.5) + i * Inches(4.25); by = Inches(1.9)
    add_rect(s, bx, by, Inches(4.0), Inches(5.1),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_text(s, icon, bx + Inches(0.15), by + Inches(0.15), Inches(0.8), Inches(0.8),
             font_size=Pt(30))
    add_text(s, title, bx + Inches(0.15), by + Inches(1.0), Inches(3.7), Inches(0.8),
             font_size=Pt(16), bold=True, color=col)
    for j, pt in enumerate(points):
        add_text(s, "▸  " + pt, bx + Inches(0.2), by + Inches(1.85) + j * Inches(0.72),
                 Inches(3.6), Inches(0.65), font_size=Pt(12), color=C_WHITE)

# ══════════════════════════════════════════════════════
# SLIDE 6 – TEAM STRUCTURE
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "02  PROJECT MANAGEMENT", Inches(0.5), Inches(0.3), fill=C_MELON_Y, text_color=C_BLACK)
add_text(s, "Cấu trúc Nhóm & Phân công", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_Y)

members = [
    ("HH", "Nguyễn Huy Hoàng", "23122031", "Project Manager\nTest Admin", C_MELON_G),
    ("NB", "Nguyễn Bá Nam",    "23122043", "Business Analyst\nQA Lead",    C_MELON_Y),
    ("TM", "Trần Tạ Quang Minh","23122042","AI/ML Engineer\nModel Evaluation", C_MELON_O),
    ("ĐK", "Nguyễn Đăng Khôi", "23122037", "Software Architect\nSAD Documentation", C_MELON_P),
    ("LV", "Lâm Hoàng Vũ",     "23122056", "Frontend Developer\nUML & DevOps", C_MELON_B),
]
# 5 columns
for i, (abbr, name, sid, role, col) in enumerate(members):
    bx = Inches(0.35) + i * Inches(2.55); by = Inches(1.9)
    bw = Inches(2.3); bh = Inches(4.8)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    # Avatar circle (fake)
    av = s.shapes.add_shape(9, bx + Inches(0.5), by + Inches(0.2), Inches(1.3), Inches(1.3))
    av.fill.solid(); av.fill.fore_color.rgb = col
    av.line.fill.background()
    add_text(s, abbr, bx + Inches(0.5), by + Inches(0.35), Inches(1.3), Inches(0.9),
             font_size=Pt(20), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    add_text(s, name, bx + Inches(0.1), by + Inches(1.65), bw - Inches(0.2), Inches(0.65),
             font_size=Pt(11), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sid,  bx + Inches(0.1), by + Inches(2.3), bw - Inches(0.2), Inches(0.4),
             font_size=Pt(10), color=C_GRAY, align=PP_ALIGN.CENTER)
    add_text(s, role, bx + Inches(0.1), by + Inches(2.75), bw - Inches(0.2), Inches(1.6),
             font_size=Pt(11), color=col, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# SLIDE 7 – PROJECT MANAGEMENT (Agile/Scrum)
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "02  PROJECT MANAGEMENT", Inches(0.5), Inches(0.3), fill=C_MELON_Y, text_color=C_BLACK)
add_text(s, "Quy trình Quản lý Dự án  ·  Agile/Scrum", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(32), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_Y)

# Sprint timeline
sprints = ["S1\nVision & Setup", "S2\nArchitecture", "S3\nBackend AI", "S4\nWorking SW", "S5\nTesting", "S6\nDemo Prep"]
for i, sp in enumerate(sprints):
    bx = Inches(0.5) + i * Inches(2.1); by = Inches(2.0)
    col = [C_MELON_G, C_MELON_Y, C_MELON_O, C_MELON_P, C_MELON_B, C_MELON_G][i]
    add_rect(s, bx, by, Inches(1.9), Inches(1.5),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    add_text(s, sp, bx + Inches(0.1), by + Inches(0.15), Inches(1.7), Inches(1.2),
             font_size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < 5:
        add_rect(s, bx + Inches(1.9), by + Inches(0.7), Inches(0.2), Inches(0.07),
                 fill_color=C_GRAY)

# Tools
tools = [
    ("🗂️ Jira", "Quản lý Sprint,\nbacklog & task tracking"),
    ("🐙 GitHub", "Version control,\nPR review, CI/CD"),
    ("💬 Discord", "Communication\n& weekly meeting"),
    ("📄 GDrive", "Tài liệu & báo cáo\ncuối Sprint"),
]
for i, (tool, desc) in enumerate(tools):
    bx = Inches(0.5) + i * Inches(3.1); by = Inches(4.0)
    add_rect(s, bx, by, Inches(2.8), Inches(2.8),
             fill_color=C_BG_CARD, border_color=C_MELON_Y, border_w=Pt(1.5))
    add_text(s, tool, bx + Inches(0.15), by + Inches(0.2), Inches(2.5), Inches(0.6),
             font_size=Pt(14), bold=True, color=C_MELON_Y)
    add_text(s, desc, bx + Inches(0.15), by + Inches(0.9), Inches(2.5), Inches(1.5),
             font_size=Pt(12), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 8 – USE-CASE MODEL
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "02  USE-CASES", Inches(0.5), Inches(0.3), fill=C_MELON_Y, text_color=C_BLACK)
add_text(s, "5 Core Use-Cases của Hệ thống", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_Y)

use_cases = [
    ("ML-AUTH", "Xác thực & Phân quyền",
     "Parent (Google OAuth) / Kid (ID+PIN) / Admin\nRole-based access control, JWT tokens",
     C_MELON_G),
    ("ML-SUBS", "Subscription & Paywall",
     "Stripe integration, Free/Pro tiers\nFeature gating dựa trên subscription plan",
     C_MELON_Y),
    ("ML-PARSE","PDF/Ảnh → Question Bank",
     "Gemini 2.5 Pro OCR + GPT-4.1 structured extraction\nToán tiếng Việt lớp 4-5, LaTeX compiler",
     C_MELON_O),
    ("ML-PRAC", "AI Tutor & Quiz (RAG)",
     "GPT-4o-mini + Pinecone Vector DB\nRAG pipeline: PDF ingestion → quiz generation",
     C_MELON_P),
    ("ML-GAMI", "Gamification System",
     "XP, Level, Badge, Leaderboard\nZustand global state, Firebase persistence",
     C_MELON_B),
]
for i, (code, title, desc, col) in enumerate(use_cases):
    row = i // 3; ci = i % 3
    bx = Inches(0.5) + ci * Inches(4.25); by = Inches(1.9) + row * Inches(2.55)
    if i == 3: bx = Inches(0.9)
    if i == 4: bx = Inches(5.15)
    bw = Inches(4.0); bh = Inches(2.3)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_rect(s, bx, by, Inches(1.1), bh, fill_color=col)
    add_text(s, code, bx + Inches(0.05), by + Inches(0.8), Inches(1.0), Inches(0.6),
             font_size=Pt(11), bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    add_text(s, title, bx + Inches(1.2), by + Inches(0.15), bw - Inches(1.3), Inches(0.65),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, desc, bx + Inches(1.2), by + Inches(0.85), bw - Inches(1.3), Inches(1.3),
             font_size=Pt(11), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 9 – NON-FUNCTIONAL REQUIREMENTS
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "03  ARCHITECTURE & NFR", Inches(0.5), Inches(0.3), fill=C_MELON_O, text_color=C_BLACK)
add_text(s, "Yêu cầu Phi chức năng (NFRs)", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_O)

nfrs = [
    ("⚡", "Performance",   "Response time < 2s\nAI streaming via SSE\nTTS cache (SHA256)", C_MELON_Y),
    ("📈", "Scalability",   "Chịu tải 1,000 CCU\nVercel edge deployment\nFirebase auto-scale",  C_MELON_G),
    ("🔐", "Security",      "JWT token auth\nSecrets trong env vars\nNext.js API Proxy layer", C_MELON_P),
    ("🛡️", "Child Safety",  "OpenAI Moderation API\n100% input/output filter\nZero-tolerance policy", C_MELON_O),
    ("🔁", "Reliability",   "CI/CD GitHub Actions\nAtomic Vercel deploy\nPR-gated merges",     C_MELON_B),
    ("🌐", "Compatibility", "Chrome/Firefox/Safari\nResponsive UI\nTailwind v4 design system", C_MELON_G),
]
for i, (icon, title, points, col) in enumerate(nfrs):
    row = i // 3; ci = i % 3
    bx = Inches(0.5) + ci * Inches(4.25); by = Inches(1.9) + row * Inches(2.55)
    bw = Inches(4.0); bh = Inches(2.3)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    add_text(s, icon + "  " + title, bx + Inches(0.2), by + Inches(0.15), bw - Inches(0.3), Inches(0.6),
             font_size=Pt(15), bold=True, color=col)
    add_text(s, points, bx + Inches(0.2), by + Inches(0.85), bw - Inches(0.3), Inches(1.3),
             font_size=Pt(12), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 10 – ARCHITECTURE OVERVIEW
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "03  ARCHITECTURE", Inches(0.5), Inches(0.3), fill=C_MELON_O, text_color=C_BLACK)
add_text(s, "Hybrid Split-Stack Architecture", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_O)

# Architecture diagram (boxes + arrows)
layers = [
    ("🌐  Browser Client", "React 19 · Next.js App Router\nTailwind v4 · Zustand · Framer Motion",
     Inches(4.5), Inches(1.85), Inches(4.5), Inches(1.1), C_MELON_B),
    ("🔷  Next.js API Proxy", "Secure Gateway · JWT Validation\nSecret Protection · Route Handlers",
     Inches(4.5), Inches(3.3), Inches(4.5), Inches(1.1), C_MELON_G),
    ("🐍  Python FastAPI Backend", "AI Engine: RAG · TTS · Vision Parsing\nPyMuPDF · LangChain · Pinecone",
     Inches(4.5), Inches(4.75), Inches(4.5), Inches(1.1), C_MELON_O),
]
for title, desc, bx, by, bw, bh, col in layers:
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_text(s, title, bx + Inches(0.15), by + Inches(0.08), bw - Inches(0.3), Inches(0.5),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, desc, bx + Inches(0.15), by + Inches(0.6), bw - Inches(0.3), Inches(0.45),
             font_size=Pt(11), color=C_GRAY)

# Arrows between layers
for ay in [Inches(2.95), Inches(4.4)]:
    add_rect(s, Inches(6.4), ay, Inches(0.5), Inches(0.35), fill_color=C_MELON_O)
    add_text(s, "↕", Inches(6.4), ay, Inches(0.5), Inches(0.35),
             font_size=Pt(14), color=C_WHITE, align=PP_ALIGN.CENTER)

# External services on right
services = [("🔥 Firebase Auth+DB", C_MELON_Y), ("📦 Pinecone VectorDB", C_MELON_P),
            ("🤖 OpenAI GPT-4o", C_MELON_G), ("💎 Gemini 2.0 Flash", C_MELON_B),
            ("🎙️ ElevenLabs TTS", C_MELON_O)]
add_text(s, "External Services", Inches(9.6), Inches(1.85), Inches(3.5), Inches(0.45),
         font_size=Pt(13), bold=True, color=C_GRAY)
for i, (svc, col) in enumerate(services):
    by = Inches(2.35) + i * Inches(0.75)
    add_rect(s, Inches(9.6), by, Inches(3.4), Inches(0.6),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(1.5))
    add_text(s, svc, Inches(9.75), by + Inches(0.1), Inches(3.1), Inches(0.45),
             font_size=Pt(12), color=col)

# ══════════════════════════════════════════════════════
# SLIDE 11 – DATA FLOW
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "03  ARCHITECTURE", Inches(0.5), Inches(0.3), fill=C_MELON_O, text_color=C_BLACK)
add_text(s, "Data Flow: PDF → Quiz (RAG Pipeline)", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(32), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_O)

steps = [
    ("1\nUpload PDF", C_MELON_B),
    ("2\nExtract Text\n(PyMuPDF)", C_MELON_G),
    ("3\nGenerate\nEmbeddings", C_MELON_Y),
    ("4\nUpsert to\nPinecone", C_MELON_P),
    ("5\nRAG Retrieval\n+ GPT-4o-mini", C_MELON_O),
    ("6\nQuiz JSON\nto Student", C_MELON_G),
]
sw = Inches(1.7); sh = Inches(1.8)
for i, (step, col) in enumerate(steps):
    bx = Inches(0.4) + i * Inches(2.15); by = Inches(2.2)
    add_rect(s, bx, by, sw, sh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    add_text(s, step, bx + Inches(0.1), by + Inches(0.2), sw - Inches(0.2), sh - Inches(0.3),
             font_size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    if i < 5:
        add_rect(s, bx + sw, by + Inches(0.75), Inches(0.45), Inches(0.08),
                 fill_color=C_GRAY)
        add_text(s, "→", bx + sw, by + Inches(0.6), Inches(0.45), Inches(0.4),
                 font_size=Pt(14), color=C_GRAY, align=PP_ALIGN.CENTER)

# TTS flow (secondary)
add_text(s, "TTS Flow: Child taps 🔊 → FastAPI → SHA256 Cache Check → ElevenLabs API → Firebase Storage → Audio URL",
         Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.55),
         font_size=Pt(13), color=C_MELON_Y)

# Moderation note
add_rect(s, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.9),
         fill_color=RGBColor(0x1A, 0x10, 0x10), border_color=C_MELON_O, border_w=Pt(1.5))
add_text(s, "🛡️  Child Safety Gate: Mọi output AI (text/audio) đều đi qua OpenAI Moderation API trước khi hiển thị cho trẻ.",
         Inches(0.7), Inches(5.2), Inches(11.8), Inches(0.65),
         font_size=Pt(13), color=C_MELON_O)

# ══════════════════════════════════════════════════════
# SLIDE 12 – TECH STACK
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "04  TECHNOLOGY", Inches(0.5), Inches(0.3), fill=C_MELON_P, text_color=C_WHITE)
add_text(s, "Technology Stack", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_P)

stacks = [
    ("🖥️  Frontend", [
        "Next.js 16 (App Router)",
        "React 19 + TypeScript",
        "Tailwind CSS v4 · shadcn/Radix",
        "Zustand · TanStack Query",
        "Framer Motion (animations)",
    ], C_MELON_B),
    ("⚙️  Backend / AI Engine", [
        "Python 3.11 + FastAPI",
        "PyMuPDF (PDF extraction)",
        "LangChain + Pinecone (RAG)",
        "OpenRouter (model routing)",
        "SSE (streaming responses)",
    ], C_MELON_G),
    ("☁️  Infrastructure & DB", [
        "Firebase Auth + Firestore",
        "Firebase Storage (assets)",
        "Pinecone (vector DB)",
        "Vercel (frontend deploy)",
        "GitHub Actions (CI/CD)",
    ], C_MELON_Y),
]
for i, (title, items, col) in enumerate(stacks):
    bx = Inches(0.5) + i * Inches(4.25); by = Inches(1.9)
    bw = Inches(4.0); bh = Inches(5.1)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    add_text(s, title, bx + Inches(0.15), by + Inches(0.15), bw - Inches(0.3), Inches(0.55),
             font_size=Pt(14), bold=True, color=col)
    for j, item in enumerate(items):
        add_text(s, "▸  " + item, bx + Inches(0.2), by + Inches(0.85) + j * Inches(0.75),
                 bw - Inches(0.3), Inches(0.65), font_size=Pt(12), color=C_WHITE)

# ══════════════════════════════════════════════════════
# SLIDE 13 – AI MODELS OVERVIEW
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "04  AI INTEGRATION", Inches(0.5), Inches(0.3), fill=C_MELON_P, text_color=C_WHITE)
add_text(s, "5 AI Models – The Brain of Melon", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_P)

models = [
    ("GPT-4o-mini", "AI Tutor Chatbot &\nRAG Quiz Generator",
     "Gia sư ảo trả lời câu hỏi theo context\nbài học, sinh MCQA quiz chính xác > 90%", C_MELON_G),
    ("Gemini 2.0/2.5", "Multimodal Vision\nProblem Parser",
     "OCR ảnh bài Toán tiếng Việt, LaTeX\nField extraction accuracy > 85%", C_MELON_B),
    ("ElevenLabs", "Text-to-Speech\n(TTS)",
     "Giọng đọc tự nhiên đa ngôn ngữ\nCaching SHA256 → giảm chi phí 70%", C_MELON_Y),
    ("Whisper-1", "Speech-to-Text\n(STT)",
     "Nhận diện giọng nói trẻ em\nVietnamese language support", C_MELON_O),
    ("OpenAI Moderation", "Content Safety\nFilter",
     "Lọc 100% output AI\nChild-safe enforcement layer", C_MELON_P),
]
for i, (name, role, desc, col) in enumerate(models):
    row = i // 3; ci = i % 3
    bx = Inches(0.5) + ci * Inches(4.25)
    if i == 3: bx = Inches(0.9)
    if i == 4: bx = Inches(5.15)
    by = Inches(1.9) + row * Inches(2.55)
    bw = Inches(4.0); bh = Inches(2.3)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_rect(s, bx, by, bw, Inches(0.08), fill_color=col)
    add_text(s, "🤖  " + name, bx + Inches(0.15), by + Inches(0.2), bw - Inches(0.3), Inches(0.55),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, role, bx + Inches(0.15), by + Inches(0.8), bw - Inches(0.3), Inches(0.55),
             font_size=Pt(11), color=C_WHITE, italic=True)
    add_text(s, desc, bx + Inches(0.15), by + Inches(1.4), bw - Inches(0.3), Inches(0.8),
             font_size=Pt(11), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 14 – EVALUATION STRATEGY
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "05  EVALUATION", Inches(0.5), Inches(0.3), fill=C_MELON_B, text_color=C_BLACK)
add_text(s, "Chiến lược Đánh giá Model AI", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_B)

add_text(s, "⚙️  Không thể fine-tune → Cải thiện qua Prompt Engineering + Evaluation Loop",
         Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.55),
         font_size=Pt(14), color=C_MELON_Y, italic=True)

evals = [
    ("1️⃣  Rubric Classifier", [
        "Task: Phân loại câu hỏi Toán theo Bloom's Taxonomy",
        "Dataset: 50 câu Toán VN lớp 4-5 (team manually labeled)",
        "Threshold: Accuracy ≥ 75%  |  Macro F1 ≥ 0.70",
        "Method: System prompt với định nghĩa Bloom's chi tiết",
    ], C_MELON_G),
    ("2️⃣  RAG Quiz Generator", [
        "Task: Sinh MCQA từ 5 lesson PDF (10 câu/doc)",
        "Dataset: 50 câu hỏi, review thủ công theo sách GK",
        "Threshold: Answer accuracy ≥ 90%  |  Distractor rate ≥ 80%",
        "Method: Strict grounding rules, no hallucination allowed",
    ], C_MELON_Y),
    ("3️⃣  PDF Problem Parser", [
        "Task: OCR 30 trang Toán VN (scan quality real-world)",
        "Baseline: Gemini 2.5 Pro ~94% DocVQA, 90% OCRBench",
        "Threshold: Field extraction ≥ 85%  |  Detection rate ≥ 90%",
        "Method: Multimodal vision + GPT-4.1 structured extraction",
    ], C_MELON_O),
]
for i, (title, points, col) in enumerate(evals):
    by = Inches(2.55) + i * Inches(1.55)
    add_rect(s, Inches(0.5), by, Inches(12.3), Inches(1.4),
             fill_color=C_BG_CARD, border_color=col, border_w=Pt(2))
    add_text(s, title, Inches(0.7), by + Inches(0.1), Inches(4), Inches(0.55),
             font_size=Pt(13), bold=True, color=col)
    line = " · ".join(points)
    add_text(s, line, Inches(0.7), by + Inches(0.65), Inches(11.8), Inches(0.65),
             font_size=Pt(11), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 15 – EXPERIMENT RESULTS
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "05  EVALUATION", Inches(0.5), Inches(0.3), fill=C_MELON_B, text_color=C_BLACK)
add_text(s, "Kết quả Thực nghiệm & Số liệu", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_B)

metrics = [
    ("Rubric Classifier",  "Accuracy",    "~78%",   "≥ 75%",  True,  C_MELON_G),
    ("Rubric Classifier",  "Macro F1",    "~0.73",  "≥ 0.70", True,  C_MELON_G),
    ("RAG Quiz Generator", "Ans Accuracy","~92%",   "≥ 90%",  True,  C_MELON_Y),
    ("RAG Quiz Generator", "Distractor",  "~85%",   "≥ 80%",  True,  C_MELON_Y),
    ("PDF Parser",         "Field Extr.", "~88%",   "≥ 85%",  True,  C_MELON_O),
    ("PDF Parser",         "Detection",   "~91%",   "≥ 90%",  True,  C_MELON_O),
]
# Header
hdrs = ["Model", "Metric", "Result", "Threshold", "Pass?"]
hw = [2.8, 2.2, 1.8, 2.2, 1.2]
hx = 0.5
for j, (hdr, w) in enumerate(zip(hdrs, hw)):
    bx = Inches(hx); by = Inches(1.85)
    add_rect(s, bx, by, Inches(w - 0.05), Inches(0.45),
             fill_color=C_MELON_B, border_color=C_BLACK, border_w=Pt(1))
    add_text(s, hdr, bx, by, Inches(w - 0.05), Inches(0.45),
             font_size=Pt(12), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    hx += w

for i, (model, metric, result, thresh, passed, col) in enumerate(metrics):
    by = Inches(2.38) + i * Inches(0.72)
    vals = [model, metric, result, thresh, "✅ PASS" if passed else "❌ FAIL"]
    ws   = [2.8, 2.2, 1.8, 2.2, 1.2]
    vx = 0.5
    for j, (v, w) in enumerate(zip(vals, ws)):
        fc = C_BG_CARD
        if j == 4: fc = RGBColor(0x08, 0x2B, 0x10) if passed else RGBColor(0x2B, 0x08, 0x08)
        add_rect(s, Inches(vx), by, Inches(w - 0.05), Inches(0.62),
                 fill_color=fc, border_color=col if j == 0 else C_BG_CARD, border_w=Pt(1.5))
        tc = col if j == 0 else (C_MELON_G if (j == 4 and passed) else C_WHITE)
        add_text(s, v, Inches(vx), by, Inches(w - 0.05), Inches(0.62),
                 font_size=Pt(12), color=tc, align=PP_ALIGN.CENTER,
                 bold=(j == 0 or j == 4))
        vx += w

add_text(s, "🔬 Key Finding: Prompt Engineering + Context grounding là yếu tố quyết định độ chính xác AI – không cần fine-tune",
         Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.55),
         font_size=Pt(12), color=C_MELON_Y, italic=True)

# ══════════════════════════════════════════════════════
# SLIDE 16 – LATENCY & SAFETY RESULTS
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "05  EVALUATION", Inches(0.5), Inches(0.3), fill=C_MELON_B, text_color=C_BLACK)
add_text(s, "Hiệu năng & Bộ lọc An toàn", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_B)

# Latency before/after
add_rect(s, Inches(0.5), Inches(1.9), Inches(6.0), Inches(4.8),
         fill_color=C_BG_CARD, border_color=C_MELON_Y, border_w=Pt(2))
add_text(s, "⚡  Latency Optimization", Inches(0.7), Inches(2.0), Inches(5.5), Inches(0.55),
         font_size=Pt(16), bold=True, color=C_MELON_Y)
perf = [
    ("❌ Before (Naive)", "RAG → LLM → TTS → Response\nTổng thời gian: ~8-12 giây", C_MELON_O),
    ("✅ After (SSE Stream)", "SSE text stream về ngay → TTS chạy ngầm\nPerceived latency: < 2 giây", C_MELON_G),
]
for i, (label, desc, col) in enumerate(perf):
    by = Inches(2.7) + i * Inches(1.7)
    add_rect(s, Inches(0.7), by, Inches(5.6), Inches(1.5),
             fill_color=C_BG_DARK, border_color=col, border_w=Pt(1.5))
    add_text(s, label, Inches(0.9), by + Inches(0.1), Inches(5.3), Inches(0.5),
             font_size=Pt(13), bold=True, color=col)
    add_text(s, desc, Inches(0.9), by + Inches(0.65), Inches(5.3), Inches(0.75),
             font_size=Pt(11), color=C_GRAY)

# Safety results
add_rect(s, Inches(6.8), Inches(1.9), Inches(6.0), Inches(4.8),
         fill_color=C_BG_CARD, border_color=C_MELON_O, border_w=Pt(2))
add_text(s, "🛡️  Safety Moderation Tests", Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.55),
         font_size=Pt(16), bold=True, color=C_MELON_O)
safety_data = [
    ("Normal content", "50 samples", "Pass 100%", C_MELON_G),
    ("Jailbreak prompts", "20 samples", "Blocked 100%", C_MELON_G),
    ("Offensive words", "30 samples", "Filtered 100%", C_MELON_G),
    ("Harmful imagery", "10 samples", "Rejected 100%", C_MELON_G),
]
for i, (test, count, result, col) in enumerate(safety_data):
    by = Inches(2.65) + i * Inches(0.85)
    add_rect(s, Inches(7.0), by, Inches(5.6), Inches(0.7),
             fill_color=C_BG_DARK, border_color=col, border_w=Pt(1))
    add_text(s, test + "  —  " + count, Inches(7.15), by + Inches(0.08), Inches(3.5), Inches(0.55),
             font_size=Pt(12), color=C_WHITE)
    add_text(s, result, Inches(10.3), by + Inches(0.08), Inches(2.0), Inches(0.55),
             font_size=Pt(12), bold=True, color=col, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════
# SLIDE 17 – DEMO SCENARIO
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "06  DEMO", Inches(0.5), Inches(0.3), fill=C_MELON_G, text_color=C_BLACK)
add_text(s, "Demo  ·  Key Scenario", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(36), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_G)

steps_demo = [
    ("Step 1", "Parent đăng nhập\n& Nạp thẻ Pro",
     "Google OAuth → Parent Dashboard\nStripe checkout → Subscription activated", C_MELON_B),
    ("Step 2", "Upload PDF Toán\n& AI Parse",
     "Gemini 2.5 Pro OCR ảnh/PDF\nGPT-4.1 extract structured JSON questions", C_MELON_G),
    ("Step 3", "Admin Review\n& Lưu Question Bank",
     "Admin xem kết quả parse\nApprove → Lưu vào Firestore + Pinecone", C_MELON_Y),
    ("Step 4", "Kid Login\n& Học tập",
     "Kid ID+PIN → Dashboard\nChọn bài → làm quiz → ElevenLabs TTS đọc đề", C_MELON_O),
    ("Step 5", "🎉 Gamification",
     "Làm đúng → XP + Pháo hoa\nNhận Badge → Leo Leaderboard", C_MELON_P),
]
for i, (num, title, desc, col) in enumerate(steps_demo):
    row = i // 3; ci = i % 3
    bx = Inches(0.5) + ci * Inches(4.25); by = Inches(1.9) + row * Inches(2.55)
    if i == 3: bx = Inches(2.6)
    if i == 4: bx = Inches(6.85)
    bw = Inches(4.0); bh = Inches(2.3)
    add_rect(s, bx, by, bw, bh, fill_color=C_BG_CARD, border_color=col, border_w=Pt(2.5))
    add_rect(s, bx, by, bw, Inches(0.45), fill_color=col)
    add_text(s, num, bx + Inches(0.15), by + Inches(0.06), bw - Inches(0.3), Inches(0.35),
             font_size=Pt(11), bold=True, color=C_BLACK)
    add_text(s, title, bx + Inches(0.15), by + Inches(0.55), bw - Inches(0.3), Inches(0.6),
             font_size=Pt(14), bold=True, color=col)
    add_text(s, desc, bx + Inches(0.15), by + Inches(1.2), bw - Inches(0.3), Inches(1.0),
             font_size=Pt(11), color=C_GRAY)

# ══════════════════════════════════════════════════════
# SLIDE 18 – LIVE DEMO (placeholder)
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "06  LIVE DEMO", Inches(0.5), Inches(0.3), fill=C_MELON_G, text_color=C_BLACK)
add_text(s, "🍉  Live Demo", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(48), bold=True, color=C_MELON_G)

add_text(s, "http://localhost:3000", Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.9),
         font_size=Pt(28), color=C_MELON_Y, align=PP_ALIGN.CENTER)

add_rect(s, Inches(1.5), Inches(3.1), Inches(10.3), Inches(3.5),
         fill_color=C_BG_CARD, border_color=C_MELON_G, border_w=Pt(3))
add_text(s, "[ Màn hình trình chiếu ứng dụng thực tế ]\n\nDemonstrate:\n▸  Kid login (em1 / 123123)\n▸  PDF Upload & AI Question Parse\n▸  Quiz session + TTS audio\n▸  Badge & XP reward",
         Inches(2.5), Inches(3.3), Inches(8.5), Inches(3.0),
         font_size=Pt(16), color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════
# SLIDE 19 – LIMITATIONS & FUTURE WORK
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

add_label_badge(s, "06  CONCLUSION", Inches(0.5), Inches(0.3), fill=C_MELON_G, text_color=C_BLACK)
add_text(s, "Hạn chế & Hướng phát triển", Inches(0.5), Inches(0.75), Inches(12), Inches(0.9),
         font_size=Pt(34), bold=True, color=C_WHITE)
add_rect(s, Inches(0.5), Inches(1.65), Inches(12.3), Inches(0.04), fill_color=C_MELON_G)

# Limitations
add_rect(s, Inches(0.5), Inches(1.9), Inches(5.8), Inches(5.1),
         fill_color=C_BG_CARD, border_color=C_MELON_O, border_w=Pt(2))
add_text(s, "⚠️  Current Limitations", Inches(0.7), Inches(2.0), Inches(5.4), Inches(0.55),
         font_size=Pt(16), bold=True, color=C_MELON_O)
limits = [
    "Parser tối ưu cho Toán VN lớp 4-5,\nchưa hỗ trợ các môn khác",
    "Chưa có LLM-based personalized\nquestion generation cho từng trẻ",
    "OCR đôi khi sai đánh số câu\nkhi ảnh chụp không rõ nét",
    "Chưa có Mobile app (chỉ Web)",
]
for i, lim in enumerate(limits):
    add_text(s, "• " + lim, Inches(0.7), Inches(2.7) + i * Inches(1.0),
             Inches(5.4), Inches(0.9), font_size=Pt(12), color=C_GRAY)

# Future
add_rect(s, Inches(6.6), Inches(1.9), Inches(6.2), Inches(5.1),
         fill_color=C_BG_CARD, border_color=C_MELON_G, border_w=Pt(2))
add_text(s, "🚀  Future Roadmap", Inches(6.8), Inches(2.0), Inches(5.8), Inches(0.55),
         font_size=Pt(16), bold=True, color=C_MELON_G)
futures = [
    "Mở rộng sang các môn Văn, Anh, Khoa học",
    "React Native mobile app (iOS + Android)",
    "Personalized AI lộ trình học tập\ndựa trên Adaptive Learning Engine",
    "Tích hợp video bài giảng + AI summary",
    "Thương mại hóa: Partnership với\ntrường học & trung tâm giáo dục",
]
for i, fut in enumerate(futures):
    add_text(s, "▸  " + fut, Inches(6.8), Inches(2.7) + i * Inches(0.88),
             Inches(5.8), Inches(0.8), font_size=Pt(12), color=C_WHITE)

# ══════════════════════════════════════════════════════
# SLIDE 20 – CONCLUSION & THANK YOU
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_bg(s, C_BG_DARK)

# Big decorative
c = s.shapes.add_shape(9, Inches(8), Inches(-1), Inches(6), Inches(6))
c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x15)
c.line.fill.background()

add_text(s, "🍉 MELON AI", Inches(1), Inches(0.6), Inches(8), Inches(0.75),
         font_size=Pt(22), bold=True, color=C_MELON_G, font_name="Calibri")

add_text(s, "Cảm ơn thầy và các bạn!", Inches(1), Inches(1.4), Inches(10), Inches(1.2),
         font_size=Pt(44), bold=True, color=C_WHITE, font_name="Calibri")

add_rect(s, Inches(1), Inches(2.7), Inches(5), Inches(0.04), fill_color=C_MELON_G)

summary = [
    "✅  Hybrid Split-Stack Architecture (Next.js + FastAPI)",
    "✅  5 AI Models phối hợp: GPT-4o, Gemini, ElevenLabs...",
    "✅  RAG Pipeline + Prompt Engineering (accuracy > 90%)",
    "✅  Child Safety Moderation – 100% filter coverage",
    "✅  Full Gamification: XP, Badge, Leaderboard",
]
for i, s_line in enumerate(summary):
    add_text(s, s_line, Inches(1), Inches(2.95) + i * Inches(0.65),
             Inches(10), Inches(0.62), font_size=Pt(14), color=C_WHITE, font_name="Calibri")

add_text(s, "Q & A", Inches(1), Inches(6.2), Inches(4), Inches(0.75),
         font_size=Pt(32), bold=True, color=C_MELON_G, font_name="Calibri")
add_text(s, "Nhóm 2  ·  SE4AI  ·  10/7/2026", Inches(1), Inches(6.9), Inches(8), Inches(0.45),
         font_size=Pt(14), color=C_GRAY, font_name="Calibri")

# ─── SAVE ─────────────────────────────────────────────
import os
out_path = r"d:\HCMUS\HK2-N3\SE4AI\melon-ai-learning-app\pa\pa5\Melon_AI_Presentation.pptx"
prs.save(out_path)
print(f"✅  Saved: {out_path}")
print(f"   Total slides: {len(prs.slides)}")
