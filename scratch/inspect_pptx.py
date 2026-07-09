import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import re

sys.stdout.reconfigure(encoding='utf-8')

prs = Presentation(r'd:\HCMUS\HK2-N3\SE4AI\melon-ai-learning-app\[Presentation] Melon.pptx')

print(f'Slide size: {prs.slide_width/914400:.2f}" x {prs.slide_height/914400:.2f}"')

# Check background and colors of original slides (1-16)
for si in [0, 1, 14, 15]:
    slide = prs.slides[si]
    xml_str = slide._element.xml
    colors = re.findall(r'srgbClr val="([0-9A-Fa-f]{6})"', xml_str[:5000])
    print(f'\nSlide {si+1} first colors: {colors[:10]}')
    
    # Check background
    bg = slide.background
    fill = bg.fill
    print(f'  bg fill type: {fill.type}')
    try: print(f'  bg fore_color: {fill.fore_color.rgb}')
    except: pass

# Check shape positions in slide 3 (an original slide)
print('\n--- Slide 3 shape positions (in inches) ---')
slide3 = prs.slides[2]
for shape in slide3.shapes:
    try:
        print(f'  {shape.name}: l={shape.left/914400:.2f}", t={shape.top/914400:.2f}", w={shape.width/914400:.2f}", h={shape.height/914400:.2f}"')
    except: pass
